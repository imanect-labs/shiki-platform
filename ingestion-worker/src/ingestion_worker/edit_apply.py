"""Office 編集 ops の適用ロジック（/edit の実体・edit.py と対）。

python-docx / openpyxl / python-pptx（全て MIT）でバイト列を開き、ops を順に適用して
バイト列へ書き戻す。ここは**純粋な変換**であり、認可・保存・バージョニングは持たない。
対象不一致は op 単位の warning に落として続行する（呼び出し元がレポートを返す）。
"""

import base64
import binascii
import io
import re
from typing import Any

# 適用結果 = (op 名, 適用数, warning|None)。DTO（edit.OpResult）への詰め替えは edit.py 側。
OpTuple = tuple[str, int, str | None]

# 埋め込み画像の敵対的入力ガード（#334・チャート静的化）。認証済みユーザー入力だが、
# サンドボックス由来入力と同様に上限を敷いて worker を守る（PIT-23 の原則）。
_MAX_IMAGE_BYTES = 8 * 1024 * 1024  # 1 枚 8MB
_MAX_IMAGES = 20  # 1 回の append で 20 枚
_IMAGE_WIDTH_INCHES = 6.0  # 本文幅に収める既定幅

# Markdown インライン記法の最小除去（**bold** / *em* / `code` → 素のテキスト）。
_INLINE_MD = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`")
_HEADING_MD = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET_MD = re.compile(r"^[-*]\s+(.*)$")
_ORDERED_MD = re.compile(r"^\d+[.)]\s+(.*)$")
# 単独行の data URL 画像（`![alt](data:image/png;base64,....)`・#334 のチャート静的化）。
# 外部 URL の画像は取得しない（worker は egress しない契約）ため data URL のみ対応する。
_IMAGE_MD = re.compile(r"^!\[([^\]]*)\]\(data:image/(png|jpeg);base64,([A-Za-z0-9+/=]+)\)$")


def _strip_inline(text: str) -> str:
    return _INLINE_MD.sub(lambda m: m.group(1) or m.group(2) or m.group(3) or "", text)


def _md_lines(markdown: str) -> list[tuple[str, str, int]]:
    """Markdown を (種別, テキスト, レベル) の行列へ落とす。

    最小集合: 見出し/箇条書き/段落＋単独行の data URL 画像（image 種別はテキスト欄に
    `alt\\x00base64` を詰める。展開は docx 側のみ・他形式は alt テキストへ縮退）。
    """
    lines: list[tuple[str, str, int]] = []
    for raw in markdown.splitlines():
        line = raw.strip()
        if not line:
            continue
        if m := _IMAGE_MD.match(line):
            lines.append(("image", f"{_strip_inline(m.group(1).strip())}\x00{m.group(3)}", 0))
        elif m := _HEADING_MD.match(line):
            lines.append(("heading", _strip_inline(m.group(2).strip()), len(m.group(1))))
        elif m := _BULLET_MD.match(line):
            lines.append(("bullet", _strip_inline(m.group(1).strip()), 0))
        elif m := _ORDERED_MD.match(line):
            lines.append(("ordered", _strip_inline(m.group(1).strip()), 0))
        else:
            lines.append(("paragraph", _strip_inline(line), 0))
    return lines


# ---------------------------------------------------------------------------
# docx
# ---------------------------------------------------------------------------

# 見出し/リストのスタイル定義を持たない docx で平文へ縮退したときの警告。
_DEGRADED_WARNING = "見出し/リストのスタイル未定義のため、一部ブロックを平文で追加しました"


def _docx_all_paragraphs(doc: Any) -> list[Any]:
    """本文＋表セル内の段落を読み順で列挙する（置換対象の走査用）。"""
    paragraphs = list(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                paragraphs.extend(cell.paragraphs)
    return paragraphs


def _replace_in_paragraph(paragraph: Any, find: str, replace: str) -> tuple[int, bool]:
    """段落内置換。run 内で置けない（run 跨ぎ）場合は段落を単一 run に潰す。

    返り値: (置換数, 書式を平坦化したか)。
    """
    count = 0
    for run in paragraph.runs:
        if find in run.text:
            count += run.text.count(find)
            run.text = run.text.replace(find, replace)
    if find not in paragraph.text:
        return count, False
    # run 跨ぎの一致: テキストを結合して置換し、先頭 run に集約（書式は先頭 run 継承）。
    remaining = paragraph.text.count(find)
    text = paragraph.text.replace(find, replace)
    for run in paragraph.runs:
        run.text = ""
    target = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    target.text = text
    return count + remaining, True


def _docx_add_heading_safe(doc: Any, text: str, level: int) -> tuple[Any, bool]:
    """見出し追加。Heading スタイル未定義の docx では太字段落へ縮退する（fail-soft）。"""
    try:
        return doc.add_heading(text, level=min(level, 9)), False
    except KeyError:
        paragraph = doc.add_paragraph()
        paragraph.add_run(text).bold = True
        return paragraph, True


def _docx_add_list_safe(doc: Any, text: str, style: str, marker: str) -> tuple[Any, bool]:
    """リスト項目追加。List スタイル未定義の docx では記号付き平文へ縮退する（fail-soft）。"""
    try:
        return doc.add_paragraph(text, style=style), False
    except KeyError:
        return doc.add_paragraph(f"{marker} {text}"), True


def _docx_add_image_safe(
    doc: Any, alt: str, b64: str, budget: list[int]
) -> tuple[Any, bool]:
    """data URL 画像を add_picture で埋め込む（#334・チャート静的化）。

    敵対的入力ガード（枚数/サイズ超過・不正 base64）に触れたら alt テキスト段落へ縮退する。
    `budget[0]` は残り枚数（呼び出し側で共有・破壊的に減算する）。
    返り値: (段落, 縮退したか)。
    """
    from docx.shared import Inches

    if budget[0] <= 0:
        return doc.add_paragraph(alt or "（図は上限枚数を超えたため省略）"), True
    try:
        data = base64.b64decode(b64, validate=True)
    except (binascii.Error, ValueError):
        return doc.add_paragraph(alt or "（図を復元できませんでした）"), True
    if not data or len(data) > _MAX_IMAGE_BYTES:
        return doc.add_paragraph(alt or "（図はサイズ上限を超えたため省略）"), True
    budget[0] -= 1
    try:
        paragraph = doc.add_paragraph()
        run = paragraph.add_run()
        run.add_picture(io.BytesIO(data), width=Inches(_IMAGE_WIDTH_INCHES))
        return paragraph, False
    except Exception:  # noqa: BLE001 - python-docx/Pillow の各種例外を alt へ縮退（fail-soft）。
        return doc.add_paragraph(alt or "（図を埋め込めませんでした）"), True


def _docx_append_blocks(doc: Any, markdown: str, anchor: Any | None = None) -> tuple[int, bool]:
    """Markdown ブロックを末尾へ追加、anchor 指定時はその直後へ移動する。

    返り値: (追加ブロック数, スタイル縮退が起きたか)。スタイル定義を持たない docx
    （最小テンプレ・他ツール生成物）でも失敗させず平文で追加する。
    """
    created = []
    degraded = False
    image_budget = [_MAX_IMAGES]  # 残り埋め込み可能枚数（_docx_add_image_safe が減算）。
    for kind, text, level in _md_lines(markdown):
        if kind == "heading":
            paragraph, fell_back = _docx_add_heading_safe(doc, text, level)
        elif kind == "bullet":
            paragraph, fell_back = _docx_add_list_safe(doc, text, "List Bullet", "•")
        elif kind == "ordered":
            paragraph, fell_back = _docx_add_list_safe(doc, text, "List Number", "-")
        elif kind == "image":
            alt, _, b64 = text.partition("\x00")
            paragraph, fell_back = _docx_add_image_safe(doc, alt, b64, image_budget)
        else:
            paragraph, fell_back = doc.add_paragraph(text), False
        created.append(paragraph)
        degraded = degraded or fell_back
    if anchor is not None:
        ref = anchor._p  # noqa: SLF001 - python-docx の要素移動は lxml 層でのみ可能。
        for paragraph in created:
            ref.addnext(paragraph._p)  # noqa: SLF001
            ref = paragraph._p  # noqa: SLF001
    return len(created), degraded


def apply_docx(data: bytes, ops: list[Any]) -> tuple[bytes, list[OpTuple]]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    results: list[OpTuple] = []
    for op in ops:
        if op.op == "replace_text":
            total, flattened = 0, False
            for paragraph in _docx_all_paragraphs(doc):
                count, flat = _replace_in_paragraph(paragraph, op.find, op.replace)
                total += count
                flattened = flattened or flat
            warning = None
            if total == 0:
                warning = f"一致なし: {op.find!r}"
            elif flattened:
                warning = "run 跨ぎの一致があり、一部段落の文字書式を平坦化しました"
            results.append((op.op, total, warning))
        elif op.op == "append_markdown":
            applied, degraded = _docx_append_blocks(doc, op.markdown)
            results.append((op.op, applied, _DEGRADED_WARNING if degraded else None))
        elif op.op == "insert_after_heading":
            anchor = next(
                (
                    p
                    for p in doc.paragraphs
                    if p.style.name.startswith("Heading") and p.text.strip() == op.heading.strip()
                ),
                None,
            )
            if anchor is None:
                results.append((op.op, 0, f"見出しが見つかりません: {op.heading!r}"))
            else:
                applied, degraded = _docx_append_blocks(doc, op.markdown, anchor)
                results.append((op.op, applied, _DEGRADED_WARNING if degraded else None))
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue(), results


# ---------------------------------------------------------------------------
# xlsx
# ---------------------------------------------------------------------------

_CELL_REF = re.compile(r"^[A-Za-z]{1,3}[1-9][0-9]{0,6}$")


def apply_xlsx(data: bytes, ops: list[Any]) -> tuple[bytes, list[OpTuple]]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data))
    results: list[OpTuple] = []
    for op in ops:
        if op.op == "add_sheet":
            if op.name in wb.sheetnames:
                results.append((op.op, 0, f"同名シートが既にあります: {op.name!r}"))
            else:
                wb.create_sheet(title=op.name)
                results.append((op.op, 1, None))
            continue
        if op.sheet not in wb.sheetnames:
            results.append((op.op, 0, f"シートが見つかりません: {op.sheet!r}"))
            continue
        ws = wb[op.sheet]
        if op.op == "set_cells":
            applied = 0
            bad_refs: list[str] = []
            for ref, value in op.cells.items():
                if not _CELL_REF.match(ref):
                    bad_refs.append(ref)
                    continue
                ws[ref.upper()] = value
                applied += 1
            warning = f"不正なセル参照を無視: {bad_refs}" if bad_refs else None
            results.append((op.op, applied, warning))
        elif op.op == "insert_rows":
            ws.insert_rows(op.at, amount=len(op.rows))
            for r_offset, row in enumerate(op.rows):
                for c_offset, value in enumerate(row):
                    ws.cell(row=op.at + r_offset, column=c_offset + 1, value=value)
            results.append((op.op, len(op.rows), None))
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue(), results


# ---------------------------------------------------------------------------
# pptx
# ---------------------------------------------------------------------------


def _pptx_paragraphs(prs: Any) -> list[Any]:
    paragraphs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs.extend(shape.text_frame.paragraphs)
    return paragraphs


def _pick_layout(prs: Any) -> Any:
    """タイトル＋本文プレースホルダを持つレイアウトを優先して選ぶ。"""
    for layout in prs.slide_layouts:
        types = {ph.placeholder_format.idx for ph in layout.placeholders}
        if 0 in types and 1 in types:
            return layout
    return prs.slide_layouts[0]


def apply_pptx(data: bytes, ops: list[Any]) -> tuple[bytes, list[OpTuple]]:
    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(io.BytesIO(data))
    results: list[OpTuple] = []
    for op in ops:
        if op.op == "replace_text":
            total, flattened = 0, False
            for paragraph in _pptx_paragraphs(prs):
                count, flat = _replace_in_paragraph(paragraph, op.find, op.replace)
                total += count
                flattened = flattened or flat
            warning = None
            if total == 0:
                warning = f"一致なし: {op.find!r}"
            elif flattened:
                warning = "run 跨ぎの一致があり、一部段落の文字書式を平坦化しました"
            results.append((op.op, total, warning))
        elif op.op == "add_slide":
            slide = prs.slides.add_slide(_pick_layout(prs))
            if slide.shapes.title is not None:
                slide.shapes.title.text = op.title
            body = next((ph for ph in slide.placeholders if ph.placeholder_format.idx == 1), None)
            if body is not None and op.bullets:
                tf = body.text_frame
                tf.text = op.bullets[0]
                for bullet in op.bullets[1:]:
                    tf.add_paragraph().text = bullet
            warning = None
            if op.bullets and body is None:
                warning = "本文プレースホルダが無いレイアウトのため bullets を配置できません"
            results.append((op.op, 1, warning))
        elif op.op == "remove_slide":
            slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001 - python-pptx に公開 API が無い定石操作。
            if op.index >= len(slide_ids):
                results.append(
                    (op.op, 0, f"スライド index が範囲外です: {op.index}（全 {len(slide_ids)} 枚）")
                )
                continue
            target = slide_ids[op.index]
            prs.part.drop_rel(target.get(qn("r:id")))
            prs.slides._sldIdLst.remove(target)  # noqa: SLF001
            results.append((op.op, 1, None))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), results
