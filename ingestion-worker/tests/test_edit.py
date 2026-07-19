"""/edit（Task 11.8）の golden テスト。

fixtures は python-docx/openpyxl/python-pptx でその場生成し、/edit を通した結果を
同ライブラリで読み戻して往復検証する（バイナリ golden ファイルを持たない）。
"""

import base64
import io

from fastapi.testclient import TestClient

DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _b64(data: bytes) -> str:
    return base64.b64encode(data).decode("ascii")


def _make_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("概要", level=1)
    doc.add_paragraph("旧製品の説明。")
    doc.add_heading("結論", level=1)
    doc.add_paragraph("旧製品を推奨する。")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _make_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "売上"
    ws["A1"] = "月"
    ws["B1"] = "金額"
    ws["A2"] = "4月"
    ws["B2"] = 100
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "旧タイトル"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _edit(client: TestClient, content_type: str, data: bytes, ops: list[dict]) -> dict:
    resp = client.post(
        "/edit",
        json={
            "tenant_id": "a-corp",
            "content_type": content_type,
            "file_name": "doc",
            "data_base64": _b64(data),
            "ops": ops,
        },
    )
    assert resp.status_code == 200, resp.text
    return resp.json()


def test_docx_replace_insert_append_roundtrip(client: TestClient) -> None:
    from docx import Document

    body = _edit(
        client,
        DOCX,
        _make_docx(),
        [
            {"op": "replace_text", "find": "旧製品", "replace": "新製品"},
            {
                "op": "insert_after_heading",
                "heading": "概要",
                "markdown": "- 追加ポイント",
            },
            {"op": "append_markdown", "markdown": "## 付録\n本文です。"},
        ],
    )
    assert body["report"]["applied_ops"] == 3
    # replace は 2 箇所（説明・結論）に効いている。
    replace = next(r for r in body["report"]["results"] if r["op"] == "replace_text")
    assert replace["applied"] == 2

    doc = Document(io.BytesIO(base64.b64decode(body["data_base64"])))
    texts = [p.text for p in doc.paragraphs]
    assert "新製品の説明。" in texts and "新製品を推奨する。" in texts
    assert not any("旧製品" in t for t in texts)
    # 見出し直後への挿入位置（「概要」の次）。
    idx = texts.index("概要")
    assert texts[idx + 1] == "追加ポイント"
    # 末尾追記（見出し＋段落）。
    assert texts[-2:] == ["付録", "本文です。"]
    heading = next(p for p in doc.paragraphs if p.text == "付録")
    assert heading.style.name.startswith("Heading")


def test_docx_append_without_styles_degrades_to_plain_text(client: TestClient) -> None:
    """Heading/List スタイル未定義の最小 docx でも append_markdown が失敗しない（fail-soft・#332）。

    他ツール生成の最小 docx（styles.xml に Heading 1 等が無い）に対して、見出しは太字段落・
    リストは記号付き平文へ縮退し、warning で縮退を報告する。
    """
    from docx import Document

    # 最小 docx: 既定テンプレから Heading/List スタイル定義を落とす（lxml 層で除去）。
    doc = Document()
    styles_el = doc.styles.element
    for style in list(styles_el):
        name_el = style.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}name")
        name = (
            name_el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
            if name_el is not None
            else ""
        )
        if name.startswith("Heading") or name in ("List Bullet", "List Number"):
            styles_el.remove(style)
    out = io.BytesIO()
    doc.save(out)

    body = _edit(
        client,
        DOCX,
        out.getvalue(),
        [{"op": "append_markdown", "markdown": "# 見出し\n- 項目\n1. 手順\n本文"}],
    )
    assert body["report"]["applied_ops"] == 1
    result = body["report"]["results"][0]
    assert result["applied"] == 4
    assert "平文で追加しました" in result["warning"]

    edited = Document(io.BytesIO(base64.b64decode(body["data_base64"])))
    texts = [p.text for p in edited.paragraphs]
    assert "見出し" in texts
    assert "• 項目" in texts
    assert "- 手順" in texts
    assert "本文" in texts


def _png_data_url(color: tuple[int, int, int] = (10, 120, 200)) -> str:
    """テスト用の小さな PNG を base64 data URL 化する（#334 の画像埋め込み）。"""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (48, 32), color).save(buf, format="PNG")
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")


def test_docx_append_embeds_data_url_image(client: TestClient) -> None:
    """単独行の data URL 画像を add_picture で埋め込む（#334・チャート静的化）。"""
    from docx import Document

    md = f"## 図表\n\n![売上チャート]({_png_data_url()})\n\n以上。"
    body = _edit(client, DOCX, _make_docx(), [{"op": "append_markdown", "markdown": md}])
    assert body["report"]["applied_ops"] == 1
    # 縮退警告は出ない（画像は正常に埋め込まれる）。
    assert body["report"]["results"][0]["warning"] is None

    edited = Document(io.BytesIO(base64.b64decode(body["data_base64"])))
    # 埋め込み画像（InlineShape）が 1 つ増える。
    assert len(edited.inline_shapes) == 1


def test_docx_append_over_image_count_degrades_to_alt(client: TestClient) -> None:
    """埋め込み画像の枚数上限（20）を超えた分は alt テキストへ縮退する（fail-soft・#334）。"""
    from docx import Document

    from ingestion_worker.edit_apply import _MAX_IMAGES

    url = _png_data_url()
    md = "\n\n".join(f"![図{i}]({url})" for i in range(_MAX_IMAGES + 3))
    body = _edit(client, DOCX, _make_docx(), [{"op": "append_markdown", "markdown": md}])
    assert body["report"]["applied_ops"] == 1
    assert "平文で追加しました" in body["report"]["results"][0]["warning"]

    edited = Document(io.BytesIO(base64.b64decode(body["data_base64"])))
    # 先頭 20 枚は画像、超過分は alt テキスト段落へ縮退する。
    assert len(edited.inline_shapes) == _MAX_IMAGES
    assert any(f"図{_MAX_IMAGES + 1}" in p.text for p in edited.paragraphs)


def test_docx_append_invalid_data_url_degrades_to_alt(client: TestClient) -> None:
    """壊れた base64 の画像は alt テキストへ縮退し、失敗させない（fail-soft・#334）。"""
    from docx import Document

    md = "![壊れた図](data:image/png;base64,@@@notbase64@@@)"
    body = _edit(client, DOCX, _make_docx(), [{"op": "append_markdown", "markdown": md}])
    # `@` は data URL 正規表現の base64 集合外なので画像行として認識されず段落になる。
    edited = Document(io.BytesIO(base64.b64decode(body["data_base64"])))
    assert len(edited.inline_shapes) == 0
    assert any("壊れた図" in p.text for p in edited.paragraphs)


def test_docx_missing_heading_is_warning_not_error(client: TestClient) -> None:
    body = _edit(
        client,
        DOCX,
        _make_docx(),
        [{"op": "insert_after_heading", "heading": "存在しない", "markdown": "x"}],
    )
    assert body["report"]["applied_ops"] == 0
    assert "見出しが見つかりません" in body["report"]["results"][0]["warning"]


def test_xlsx_set_cells_insert_rows_add_sheet(client: TestClient) -> None:
    from openpyxl import load_workbook

    body = _edit(
        client,
        XLSX,
        _make_xlsx(),
        [
            {"op": "set_cells", "sheet": "売上", "cells": {"B2": 250, "C1": "備考"}},
            {"op": "insert_rows", "sheet": "売上", "at": 3, "rows": [["5月", 300]]},
            {"op": "add_sheet", "name": "集計"},
        ],
    )
    assert body["report"]["applied_ops"] == 3
    wb = load_workbook(io.BytesIO(base64.b64decode(body["data_base64"])))
    ws = wb["売上"]
    assert ws["B2"].value == 250
    assert ws["C1"].value == "備考"
    assert ws["A3"].value == "5月" and ws["B3"].value == 300
    assert "集計" in wb.sheetnames


def test_xlsx_unknown_sheet_is_warning(client: TestClient) -> None:
    body = _edit(
        client,
        XLSX,
        _make_xlsx(),
        [{"op": "set_cells", "sheet": "無い", "cells": {"A1": 1}}],
    )
    assert body["report"]["applied_ops"] == 0
    assert "シートが見つかりません" in body["report"]["results"][0]["warning"]


def test_pptx_replace_add_remove_roundtrip(client: TestClient) -> None:
    from pptx import Presentation

    body = _edit(
        client,
        PPTX,
        _make_pptx(),
        [
            {"op": "replace_text", "find": "旧タイトル", "replace": "新タイトル"},
            {"op": "add_slide", "title": "まとめ", "bullets": ["一点目", "二点目"]},
        ],
    )
    assert body["report"]["applied_ops"] == 2
    prs = Presentation(io.BytesIO(base64.b64decode(body["data_base64"])))
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "新タイトル"
    texts = [
        p.text
        for shape in prs.slides[1].shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
    ]
    assert "まとめ" in texts and "一点目" in texts and "二点目" in texts

    # 追加したスライドを削除して 1 枚へ戻す（往復）。
    body = _edit(
        client,
        PPTX,
        base64.b64decode(body["data_base64"]),
        [{"op": "remove_slide", "index": 1}],
    )
    prs = Presentation(io.BytesIO(base64.b64decode(body["data_base64"])))
    assert len(prs.slides) == 1


def test_unknown_op_is_422(client: TestClient) -> None:
    resp = client.post(
        "/edit",
        json={
            "tenant_id": "a",
            "content_type": DOCX,
            "file_name": "d",
            "data_base64": _b64(_make_docx()),
            "ops": [{"op": "drop_everything"}],
        },
    )
    assert resp.status_code == 422
    assert resp.json()["detail"]["error"] == "invalid_ops"


def test_kind_mismatched_op_is_422(client: TestClient) -> None:
    # xlsx 用 op を docx へ投げても受理しない（種別ごとのクローズド集合）。
    resp = client.post(
        "/edit",
        json={
            "tenant_id": "a",
            "content_type": DOCX,
            "file_name": "d",
            "data_base64": _b64(_make_docx()),
            "ops": [{"op": "set_cells", "sheet": "S", "cells": {"A1": 1}}],
        },
    )
    assert resp.status_code == 422


def test_unsupported_content_type_is_422(client: TestClient) -> None:
    resp = client.post(
        "/edit",
        json={
            "tenant_id": "a",
            "content_type": "application/pdf",
            "file_name": "d",
            "data_base64": _b64(b"x"),
            "ops": [{"op": "replace_text", "find": "a", "replace": "b"}],
        },
    )
    assert resp.status_code == 422
    assert resp.json()["detail"]["error"] == "unsupported_content_type"


def test_broken_file_is_422(client: TestClient) -> None:
    resp = client.post(
        "/edit",
        json={
            "tenant_id": "a",
            "content_type": DOCX,
            "file_name": "d",
            "data_base64": _b64(b"not a docx"),
            "ops": [{"op": "replace_text", "find": "a", "replace": "b"}],
        },
    )
    assert resp.status_code == 422
    assert resp.json()["detail"]["error"] == "edit_failed"


def test_size_limit_is_enforced(client: TestClient, monkeypatch) -> None:
    from ingestion_worker.settings import get_settings

    settings = get_settings()
    monkeypatch.setattr(settings, "max_edit_bytes", 10)
    try:
        resp = client.post(
            "/edit",
            json={
                "tenant_id": "a",
                "content_type": DOCX,
                "file_name": "d",
                "data_base64": _b64(_make_docx()),
                "ops": [{"op": "replace_text", "find": "a", "replace": "b"}],
            },
        )
        assert resp.status_code == 422
        assert resp.json()["detail"]["error"] == "too_large"
    finally:
        get_settings.cache_clear()
